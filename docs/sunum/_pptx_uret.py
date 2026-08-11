# -*- coding: utf-8 -*-
"""RemaLab WMS sunumunu .pptx olarak üretir.

Tasarım dili HTML sunumla aynı: paleti uygulamanın kendi token'larından alır,
statü/onarım kodlarını yapısal öğe olarak kullanır, rakamları mono yazı tipiyle
dizer. 16:9, PowerPoint'te düzenlenebilir gerçek şekiller (resim değil).

Çalıştırma:  python docs/sunum/_pptx_uret.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

KOK = os.path.dirname(os.path.abspath(__file__))
EKRAN = os.path.join(KOK, "ekran_goruntuleri")
CIKTI = os.path.join(KOK, "RemaLab_WMS_Sunum.pptx")

# ── Palet: RemaLab arayüzünün kendi değerleri ───────────────────────────────
INK      = RGBColor(0x0F, 0x11, 0x18)   # zemin
PANEL    = RGBColor(0x17, 0x1A, 0x24)   # kart
PANEL2   = RGBColor(0x1E, 0x22, 0x30)
LINE     = RGBColor(0x27, 0x2B, 0x39)
TEXT     = RGBColor(0xED, 0xEF, 0xF5)
MUTED    = RGBColor(0x88, 0x92, 0xB5)   # arayüzde birebir bu renk
SIGNAL   = RGBColor(0xE9, 0xA1, 0x3B)   # sinyal amberi — yapısal aksan
DONE     = RGBColor(0x34, 0xD3, 0x99)
WAIT     = RGBColor(0xA7, 0x8B, 0xFA)
STOP     = RGBColor(0xF8, 0x71, 0x71)
LIVE     = RGBColor(0x60, 0xA5, 0xFA)

SANS = "Segoe UI"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)      # 16:9
PAD = Inches(0.72)
ICER = W - 2 * PAD                       # içerik genişliği


def yeni_sunum():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def slayt(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])          # boş düzen
    zemin = s.shapes.add_shape(1, 0, 0, W, H)               # dikdörtgen
    zemin.fill.solid(); zemin.fill.fore_color.rgb = INK
    zemin.line.fill.background(); zemin.shadow.inherit = False
    return s


def kutu(s, x, y, w, h, dolgu=None, kenar=None, kalinlik=0.75, radius=None):
    """Kart / şerit. radius verilirse yuvarlatılmış dikdörtgen (şekil 5)."""
    sekil = s.shapes.add_shape(5 if radius else 1, x, y, w, h)
    if radius:
        sekil.adjustments[0] = radius
    if dolgu is not None:
        sekil.fill.solid(); sekil.fill.fore_color.rgb = dolgu
    else:
        sekil.fill.background()
    if kenar is not None:
        sekil.line.color.rgb = kenar; sekil.line.width = Pt(kalinlik)
    else:
        sekil.line.fill.background()
    sekil.shadow.inherit = False
    return sekil


def yaz(s, x, y, w, h, satirlar, hiza=PP_ALIGN.LEFT, capa=MSO_ANCHOR.TOP):
    """satirlar: (metin, punto, renk, kalin, font, satir_araligi, bosluk_ust) sözlükleri."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = capa
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, sat in enumerate(satirlar):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = hiza
        if sat.get("ara"):
            p.line_spacing = sat["ara"]
        if sat.get("ust"):
            p.space_before = Pt(sat["ust"])
        if sat.get("alt"):
            p.space_after = Pt(sat["alt"])
        r = p.add_run(); r.text = sat["t"]
        f = r.font
        f.size = Pt(sat.get("pt", 14))
        f.color.rgb = sat.get("c", TEXT)
        f.bold = sat.get("kalin", False)
        f.name = sat.get("font", SANS)
    return tb


def eyebrow(s, kod, bolum):
    """Üst şerit: kod + bölüm adı + amber ince çizgi."""
    yaz(s, PAD, Inches(0.52), Inches(2.4), Inches(0.3),
        [{"t": kod, "pt": 11, "c": SIGNAL, "kalin": True, "font": MONO}])
    yaz(s, PAD + Inches(2.5), Inches(0.52), ICER - Inches(2.5), Inches(0.3),
        [{"t": bolum.upper(), "pt": 11, "c": MUTED, "font": MONO}], hiza=PP_ALIGN.RIGHT)
    cizgi = kutu(s, PAD, Inches(0.86), ICER, Emu(9525), dolgu=LINE)
    return cizgi


def baslik(s, metin, pt=34, y=Inches(1.12)):
    yaz(s, PAD, y, ICER, Inches(1.0),
        [{"t": metin, "pt": pt, "c": TEXT, "kalin": True, "ara": 1.05}])


def lede(s, metin, y, pt=14.5, w=None):
    yaz(s, PAD, y, w or Inches(9.6), Inches(1.1),
        [{"t": metin, "pt": pt, "c": MUTED, "ara": 1.35}])


def kart(s, x, y, w, h, etiket, bas, govde, kenar=LINE, sol_serit=None):
    kutu(s, x, y, w, h, dolgu=PANEL, kenar=kenar, radius=0.05)
    if sol_serit:
        kutu(s, x, y, Inches(0.045), h, dolgu=sol_serit)
    ic = Inches(0.26)
    sat = []
    if etiket:
        sat.append({"t": etiket, "pt": 9.5, "c": SIGNAL, "kalin": True, "font": MONO, "alt": 6})
    if bas:
        sat.append({"t": bas, "pt": 13.5, "c": TEXT, "kalin": True, "alt": 5, "ara": 1.1})
    if govde:
        sat.append({"t": govde, "pt": 11, "c": MUTED, "ara": 1.3})
    yaz(s, x + ic, y + ic, w - 2 * ic, h - 2 * ic, sat)


def istatistik(s, x, y, w, h, deger, anahtar, renk=TEXT):
    kutu(s, x, y, w, h, dolgu=PANEL, kenar=LINE, radius=0.06)
    yaz(s, x + Inches(0.24), y + Inches(0.2), w - Inches(0.48), h - Inches(0.4),
        [{"t": deger, "pt": 33, "c": renk, "kalin": True, "font": MONO, "alt": 4},
         {"t": anahtar, "pt": 10.5, "c": MUTED, "ara": 1.25}])


def statu_kutu(s, x, y, w, h, kod, ad, renk):
    kutu(s, x, y, w, h, dolgu=PANEL, kenar=LINE, radius=0.08)
    kutu(s, x, y, w, Inches(0.032), dolgu=renk)          # üst şerit = seviye rengi
    yaz(s, x + Inches(0.14), y + Inches(0.16), w - Inches(0.28), h - Inches(0.3),
        [{"t": kod, "pt": 15, "c": renk, "kalin": True, "font": MONO, "alt": 3},
         {"t": ad, "pt": 8.5, "c": MUTED, "ara": 1.15}])


def not_serit(s, y, metin, h=Inches(0.62)):
    kutu(s, PAD, y, Inches(0.028), h, dolgu=SIGNAL)
    yaz(s, PAD + Inches(0.2), y, ICER - Inches(0.2), h,
        [{"t": metin, "pt": 11.5, "c": MUTED, "ara": 1.35}], capa=MSO_ANCHOR.MIDDLE)


def madde(s, x, y, w, h, maddeler, pt=11.5):
    sat = []
    for m in maddeler:
        sat.append({"t": "—  " + m, "pt": pt, "c": MUTED, "ara": 1.3, "alt": 7})
    yaz(s, x, y, w, h, sat)


def tablo(s, x, y, w, basliklar, satirlar, gen, pt=10.5):
    """Elle dizilmiş tablo — PowerPoint tablo nesnesi tema rengini zorluyor."""
    yh = Inches(0.34)
    cx = x
    for i, b in enumerate(basliklar):
        yaz(s, cx, y, gen[i], yh, [{"t": b.upper(), "pt": 9, "c": MUTED, "kalin": True, "font": MONO}])
        cx += gen[i]
    kutu(s, x, y + yh - Inches(0.06), w, Emu(9525), dolgu=LINE)
    sy = y + yh + Inches(0.06)
    for sat in satirlar:
        cx = x
        yuk = Inches(0.30) if len(sat[-1]) < 62 else Inches(0.44)
        for i, hucre in enumerate(sat):
            renk = sat[-1] if False else None
            f = MONO if i == 0 else SANS
            c = TEXT if i <= 1 else MUTED
            kal = i == 0
            if isinstance(hucre, tuple):
                metin, c = hucre
            else:
                metin = hucre
            yaz(s, cx, sy, gen[i] - Inches(0.1), yuk,
                [{"t": metin, "pt": pt, "c": c, "kalin": kal, "font": f, "ara": 1.2}])
            cx += gen[i]
        sy += yuk
        kutu(s, x, sy - Inches(0.05), w, Emu(9525), dolgu=PANEL2)
    return sy


def resim(s, ad, x, y, w):
    p = os.path.join(EKRAN, ad)
    pic = s.shapes.add_picture(p, x, y, width=w)
    cer = kutu(s, x, y, w, pic.height, dolgu=None, kenar=LINE)
    cer.line.width = Pt(0.75)
    return pic


# ═══════════════════════════════════════════════════════════════════════════
prs = yeni_sunum()

# ── 01 Başlık ──────────────────────────────────────────────────────────────
s = slayt(prs)
kutu(s, 0, 0, Inches(0.09), H, dolgu=SIGNAL)
yaz(s, PAD, Inches(1.55), ICER, Inches(0.4),
    [{"t": "REMALAB WMS", "pt": 12, "c": SIGNAL, "kalin": True, "font": MONO}])
yaz(s, PAD, Inches(2.05), Inches(10.8), Inches(2.2),
    [{"t": "Telefon yenileme hattını\ntek bir sistemde yönetmek",
      "pt": 44, "c": TEXT, "kalin": True, "ara": 1.02}])
lede(s, "Depo kabulünden müşteriye sevkiyata kadar cihazın geçtiği her adımı, her onarımı, "
        "her parça hareketini ve her testi kayda alan masaüstü uygulaması.",
     Inches(4.35), pt=15, w=Inches(9.2))
kutu(s, PAD, Inches(5.55), Inches(3.4), Emu(28575), dolgu=SIGNAL)
yaz(s, PAD, Inches(5.8), ICER, Inches(0.6),
    [{"t": "Depo  ·  Test  ·  Demontaj  ·  Onarım  ·  Sevkiyat",
      "pt": 12.5, "c": MUTED, "font": MONO}])

# ── 02 Problem ─────────────────────────────────────────────────────────────
s = slayt(prs); eyebrow(s, "01", "Problem")
baslik(s, "Hattın belleği kâğıtta ve kafalardaydı")
lede(s, "Bir cihaza kimin, hangi parçayı ne zaman taktığı; testten neden döndüğü; "
        "müşterinin neye onay verdiği — hepsi ayrı yerlerde tutuluyordu.", Inches(1.95),
     w=Inches(7.3))
madde(s, PAD, Inches(2.95), Inches(7.1), Inches(3.2), [
    "Cihaz hangi aşamada, kimde bekliyor belirsiz",
    "Aynı parça iki kez talep edilebiliyor, stok tutmuyor",
    "Onarım sırası kişiye göre değişiyor; L1 ekibi bitmemiş bir ekranın üstünde çalışıyor",
    "Müşteri onayı gereken iş, onay alınmadan yapılabiliyor",
    "Test sonuçları cihazın geçmişine bağlanmıyor",
], pt=12.5)
gx, gy, gw, gh = Inches(8.15), Inches(2.0), Inches(2.15), Inches(1.5)
istatistik(s, gx, gy, gw, gh, "24", "cihaz statüsü —\nkabulden sevkiyata", SIGNAL)
istatistik(s, gx + gw + Inches(0.22), gy, gw, gh, "9", "onarım statüsü —\natamadan tamamlamaya", SIGNAL)
istatistik(s, gx, gy + gh + Inches(0.22), gw, gh, "7", "onarım ekibi,\nsıralı çalışan", SIGNAL)
istatistik(s, gx + gw + Inches(0.22), gy + gh + Inches(0.22), gw, gh, "3", "test noktası,\nPhoneCheck bağlı", SIGNAL)

# ── 03 Mimari ──────────────────────────────────────────────────────────────
s = slayt(prs); eyebrow(s, "02", "Mimari")
baslik(s, "Tek pencere, iki dünya")
lede(s, "Kullanıcı tek bir masaüstü uygulaması görüyor. İçinde modern bir web arayüzü ile "
        "Python iş mantığı yan yana çalışıyor; ikisi QWebChannel köprüsüyle konuşuyor.",
     Inches(1.95), w=Inches(10.5))
kw, kh, ky = Inches(3.83), Inches(2.35), Inches(2.95)
kart(s, PAD, ky, kw, kh, "ARAYÜZ", "React 19 + Vite + Tailwind",
     "Qt WebEngine (Chromium) içinde açılan tek sayfalı uygulama. 63 dosya, ~33.500 satır. "
     "Ekranlar rol bazlı menülerle ayrılmış.")
kart(s, PAD + kw + Inches(0.26), ky, kw, kh, "KÖPRÜ", "QWebChannel",
     "Arayüzdeki her api.xxx() çağrısı Python tarafında bir @Slot'a düşer ve JSON döner. "
     "Masaüstünde doğrudan, tarayıcıda WebSocket üzerinden.", kenar=SIGNAL)
kart(s, PAD + 2 * (kw + Inches(0.26)), ky, kw, kh, "İŞ MANTIĞI + VERİ",
     "Python + PySide6 + SQLAlchemy",
     "Bütün kurallar burada. core/web_bridge.py sistemin kalbi; uzak PostgreSQL "
     "sunucusuna SQLAlchemy ile bağlanır.")
not_serit(s, Inches(5.7),
          "Neden önemli: kural arayüzde de var ama arayüz hiçbir zaman tek bekçi değil. "
          "Bir ekran atlanırsa ya da eski bir sürüm çalışırsa kural yine backend'de uygulanıyor.")

# ── 04 Statü akışı ─────────────────────────────────────────────────────────
s = slayt(prs); eyebrow(s, "03", "Cihazın yolculuğu")
baslik(s, "Statü akışı — kabulden sevkiyata", pt=31)
lede(s, "Her cihaz numaralı bir durum makinesinde ilerler. Geçişler serbest değil: hangi "
        "statüden hangisine gidilebileceği tanımlı (service_statu_map).", Inches(1.9),
     pt=13, w=Inches(11.5))
akis = [
    [("100", "Ön bildirim", MUTED), ("101", "Depo kabulü", MUTED), ("102", "İlk teste aktarıldı", LIVE),
     ("103", "İlk test bekleniyor", LIVE), ("104", "İlk test tamamlandı", LIVE), ("105", "Planlama onayı", SIGNAL)],
    [("106", "Onaya gönderilecek", WAIT), ("107", "Müşteri onayı bekleniyor", WAIT), ("136", "Onay / red alındı", WAIT),
     ("109", "Üretim aşamasında", DONE), ("138", "Ara test bekleniyor", LIVE), ("124", "Son teste teslim", LIVE)],
    [("125", "Son teste kabul", LIVE), ("126", "Depoya sevk", DONE), ("127", "Müşteriye sevkiyat", DONE),
     ("134", "RMA incelemesi", STOP), ("135", "İade öncesi test", STOP), ("128", "Serbest bırakıldı", MUTED)],
]
bw = (ICER - Inches(0.14) * 5) / 6
by = Inches(2.85)
for satir in akis:
    for i, (kod, ad, renk) in enumerate(satir):
        statu_kutu(s, PAD + i * (bw + Inches(0.14)), by, bw, Inches(0.86), kod, ad, renk)
    by += Inches(1.02)
not_serit(s, Inches(6.15),
          "Test başarısızsa cihaz geri döner: ara testten 138 → 109, son testten 125 → 109. "
          "İkisinde de L1/L2 onarımları yeniden teknisyene atanır.", h=Inches(0.55))

# ── 05 Roller ──────────────────────────────────────────────────────────────
s = slayt(prs); eyebrow(s, "04", "Roller")
baslik(s, "Herkes yalnız kendi kapısını görür")
lede(s, "Menü role göre açılır ve her statü geçişi yetkiye bağlı. Depocu onarım kapatamaz, "
        "teknisyen sevkiyat yapamaz.", Inches(1.95), w=Inches(10.5))
roller = [
    ("SPA_P", "Yedek Parça Personeli", "Kayıt kabul 100→101, ilk teste aktarma 101→102, müşteriye sevk 126→127."),
    ("QAC", "Test Personeli", "İlk teste kabul 102→103, üretime teslim 103→104, son teste kabul 124→125, sonuç 125→126/109."),
    ("TEC_DISMANTLE", "Demontaj Teknisyeni", "Teknik departmana kabul 104→105, onarım planı, Üretime Aktar kararı, onaya gönderme 105→106."),
    ("ÜRETİM", "Üretim Teknisyeni", "Üretim kaydını görüntüleme, teknisyen atama, parça teslim alma, onarımı tamamlama."),
    ("HAVUZ", "Onarım Havuzu", "Batarya · Kamera · Ekran · Kasa · L1 · L2 · L3 — her ekibin kendi iş listesi."),
    ("MNG / AS", "Ara Test & Yönetim", "Ara test 138→124, müşteri onayı, fiyat matrisleri, departman ve kullanıcı yönetimi."),
]
kw, kh = Inches(3.83), Inches(1.75)
for i, (et, bas, gov) in enumerate(roller):
    x = PAD + (i % 3) * (kw + Inches(0.26))
    y = Inches(2.85) + (i // 3) * (kh + Inches(0.24))
    kart(s, x, y, kw, kh, et, bas, gov)

# ── 06 Veri modeli ─────────────────────────────────────────────────────────
s = slayt(prs); eyebrow(s, "05", "Veri modeli")
baslik(s, "Onarım nedir? Satır değil, grup.")
lede(s, "repair_records tablosunda her satır bir PARÇADIR. Bir ONARIM ise aynı cihaz + aynı "
        "görev grubu altındaki TÜM satırlardır.", Inches(1.95), w=Inches(6.9))
madde(s, PAD, Inches(3.0), Inches(6.7), Inches(2.4), [
    "Teknisyen ataması grubun hepsini birden atar",
    "Tamamlama da grubun hepsini birden kapatır",
    "Yeniden açma tam tersi yönde çalışır: kapalı satırları açar",
    "Bir satır bile engelliyse hiçbiri işlenmez — yarım kapanmış onarım oluşamaz",
], pt=12)
not_serit(s, Inches(5.5),
          "Bu ayrım sistemin en kritik detayı. Tamamlama tek satırı kapattığı sürece cihaz "
          "\"bitti\" görünüyor ama alt seviye onarım sonsuza kadar bekliyordu.")
# örnek tablo
tx, ty, tw = Inches(7.85), Inches(2.35), Inches(4.75)
kutu(s, tx, ty, tw, Inches(3.35), dolgu=PANEL, kenar=LINE, radius=0.04)
yaz(s, tx + Inches(0.26), ty + Inches(0.24), tw - Inches(0.5), Inches(0.3),
    [{"t": "ÖRNEK — BİR CİHAZ", "pt": 9.5, "c": SIGNAL, "kalin": True, "font": MONO}])
tablo(s, tx + Inches(0.26), ty + Inches(0.66), tw - Inches(0.5),
      ["Grup", "Parça", "Durum"],
      [["DISPLAY", "LCD", ("1002", DONE)],
       ["BATTERY", "Battery", ("1002", DONE)],
       ["BATTERY", "Battery Flex", ("1002", DONE)],
       ["CASE", "Back Cover", ("1001", LIVE)],
       ["L1REPAIR", "Receiver", ("1004", WAIT)]],
      [Inches(1.55), Inches(1.7), Inches(1.0)])
yaz(s, tx + Inches(0.26), ty + Inches(2.78), tw - Inches(0.5), Inches(0.5),
    [{"t": "Batarya onarımı iki satır, ikisi birlikte kapandı. L1 bekliyor çünkü Kasa hâlâ açık.",
      "pt": 10, "c": MUTED, "ara": 1.25}])

# ── 07 Onarım statüleri ────────────────────────────────────────────────────
s = slayt(prs); eyebrow(s, "06", "Onarım statüleri")
baslik(s, "Bir onarımın geçebileceği dokuz durum", pt=31)
tablo(s, PAD, Inches(2.0), ICER,
      ["Kod", "Anlamı", "Ne zaman oluşur"],
      [[("1000", MUTED), "Teknisyene Atanacak", "Onarım açıldı, havuzda bekliyor. Henüz kimseye atanmamış."],
       [("1001", LIVE), "Teknisyen Atandı", "İş bir teknisyende. Parça teslimi ve tamamlama bu durumda yapılır."],
       [("1002", DONE), "Onarım Tamamlandı", "Bütün şartlar sağlandı ve iş kapandı. closed_at yazılır."],
       [("1003", STOP), "Onarım İptal Edildi", "İş yapılmayacak. Kural motorlarında hiç sayılmaz."],
       [("1004", WAIT), "Sırada Bekliyor", "Kendinden yüksek seviyeli onarım açık. Teknisyen atanabilir; parça teslimi ve tamamlama kapalı."],
       [("1005", SIGNAL), "Müşteri Onayı Bekliyor", "İş, müşteri onayı olmadan yapılamaz."],
       [("1006", LIVE), "Onarım Testi Bekleniyor", "Kamera / L3 / Ekran / Kasa tamamlanınca doğrudan 1002 olmaz — önce bitiş testi."],
       [("1007", SIGNAL), "Onarım Testi Başarısız", "Bitiş testinden döndü, yeniden çalışılacak."],
       [("1008", SIGNAL), "Parça Bekleniyor", "Depoda parça yok, iş parça gelene kadar duruyor."]],
      [Inches(1.0), Inches(2.9), Inches(8.0)], pt=11)

# ── 08 Seviyelendirme ──────────────────────────────────────────────────────
s = slayt(prs); eyebrow(s, "07", "Seviyelendirme")
baslik(s, "Onarımlar rastgele değil, sırayla yapılır")
lede(s, "Sıra mission_groups.order_number ile tanımlı: BÜYÜK olan önce, EŞİT olanlar paralel. "
        "Sırası gelmeyen onarım 1004 ile beklemeye alınır, üst seviye kapanınca kendiliğinden açılır.",
     Inches(1.9), pt=13, w=Inches(11.6))
seviyeler = [
    ("99", "RMA Kontrol ve Onarımı", "Her şeyden önce gelir", "1. sıra", STOP),
    ("9", "L3 Onarımı", "Anakart / mikro lehim seviyesi", "2. sıra", SIGNAL),
    ("7", "Batarya · Kamera · Kasa · Ekran", "Dördü aynı seviyede — birbirini beklemez, paralel çalışır", "3. sıra", LIVE),
    ("6", "L1 / L2 Onarımı", "En son. Bittiğinde cihaz kendiliğinden Ara Test'e geçer", "4. sıra", DONE),
    ("—", "Demontaj · Test", "Sıra numarası yok: ne bekletir ne bekletilir, her zaman görünür", "sırasız", MUTED),
]
sy = Inches(2.9)
for ordn, ad, aciklama, etiket, renk in seviyeler:
    h = Inches(0.78)
    kutu(s, PAD, sy, ICER, h, dolgu=PANEL, kenar=LINE, radius=0.09)
    kutu(s, PAD, sy, Inches(0.05), h, dolgu=renk)
    yaz(s, PAD + Inches(0.32), sy, Inches(0.9), h,
        [{"t": ordn, "pt": 21, "c": renk, "kalin": True, "font": MONO}], capa=MSO_ANCHOR.MIDDLE)
    yaz(s, PAD + Inches(1.3), sy, Inches(8.6), h,
        [{"t": ad, "pt": 14, "c": TEXT, "kalin": True, "alt": 3},
         {"t": aciklama, "pt": 10.5, "c": MUTED}], capa=MSO_ANCHOR.MIDDLE)
    yaz(s, PAD + Inches(10.1), sy, Inches(1.7), h,
        [{"t": etiket, "pt": 10.5, "c": renk, "kalin": True, "font": MONO}],
        hiza=PP_ALIGN.RIGHT, capa=MSO_ANCHOR.MIDDLE)
    sy += h + Inches(0.14)

# ── 09 Çakışma kuralları ───────────────────────────────────────────────────
s = slayt(prs); eyebrow(s, "08", "Kural motoru")
baslik(s, "Sistem yanlış parça kombinasyonuna izin vermiyor", pt=30)
lede(s, "Bazı parçalar aynı cihazda birlikte olamaz — LCD komple ekran modülüdür, ön cam ve "
        "arka aydınlatma zaten onun içinde gelir. Kural cihaz genelinde çalışır.", Inches(1.9),
     pt=13, w=Inches(11.5))
kw = Inches(5.9)
kutu(s, PAD, Inches(2.85), kw, Inches(2.75), dolgu=PANEL, kenar=LINE, radius=0.05)
yaz(s, PAD + Inches(0.26), Inches(3.08), kw - Inches(0.5), Inches(0.3),
    [{"t": "ÇAKIŞAN İKİLİLER", "pt": 9.5, "c": SIGNAL, "kalin": True, "font": MONO}])
madde(s, PAD + Inches(0.26), Inches(3.5), kw - Inches(0.5), Inches(2.0), [
    "LCD  ↔  Front Glass · Front Bezel · Backlight · Touch Glass",
    "Touch Glass  ↔  Front Glass with POL",
    "Front Glass  ↔  Front Glass with POL",
    "Main Camera  ↔  5x · 1x · 0,5x objektifler",
    "Back Cover  ↔  Middle Frame · Back Glass",
    "Front Camera / _Y  ↔  Front Camera_R",
], pt=11)
x2 = PAD + kw + Inches(0.3)
kutu(s, x2, Inches(2.85), kw, Inches(2.75), dolgu=PANEL, kenar=LINE, radius=0.05)
yaz(s, x2 + Inches(0.26), Inches(3.08), kw - Inches(0.5), Inches(0.3),
    [{"t": "\"YALNIZCA BİR TANE\" GRUPLARI", "pt": 9.5, "c": SIGNAL, "kalin": True, "font": MONO}])
madde(s, x2 + Inches(0.26), Inches(3.5), kw - Inches(0.5), Inches(1.3), [
    "Batarya flexi — Battery / Cracked / Diag / Ti arasından bir tane",
    "Kombo flexler — On-Off · Volume Key · Receiver · Flash · NFC ve \"+\" ile paketlenmiş hâlleri",
], pt=11)
yaz(s, x2 + Inches(0.26), Inches(4.85), kw - Inches(0.5), Inches(0.7),
    [{"t": "Ayrıca: aynı parça iki kez eklenemez, tamamlanmış onarıma parça eklenemez, "
           "iptal edilen parça yeniden eklenebilir.", "pt": 10.5, "c": MUTED, "ara": 1.25}])
not_serit(s, Inches(5.85),
          "Kural hem ekleme hem düzenleme yolunda uygulanıyor. Ekranda çakışan parçalar listede "
          "üstü çizili ve sebebi yazılı görünür — kullanıcı reddedilecek bir istek göndermez.")

# ── 10 Tamamlama şartları ──────────────────────────────────────────────────
s = slayt(prs); eyebrow(s, "09", "Tamamlama şartları")
baslik(s, "Bir onarım hangi durumda kapanabilir?")
lede(s, "Dört şart tek bir yerde toplandı. Aynı kuralı hem normal ekran hem hızlı bitiş "
        "ekranı çağırır — kural iki yere yazılırsa zamanla ayrışır.", Inches(1.95), w=Inches(11))
sartlar = [
    ("ŞART 1", "Teknisyen atanmış olmalı", "Kimin yaptığı bilinmeyen iş kapanamaz.", LINE),
    ("ŞART 2", "Parça depodan çıkmış olmalı", "Stok takipli parçada \"Stoktan Çıktı\" şartı aranır; takipsizde aranmaz.", LINE),
    ("ŞART 3", "Müşteri onayı alınmış olmalı", "Onay gerektiren akışlarda cihaz 106/107/136'da park ederken kapatılamaz.", LINE),
    ("ŞART 4", "Sırası gelmiş olmalı", "Kendinden yüksek seviyede açık onarım varken kapanmaz. Parçasız kayıt sayılmaz.", SIGNAL),
]
kw2 = (ICER - Inches(0.24) * 3) / 4
for i, (et, bas, gov, kenar) in enumerate(sartlar):
    kart(s, PAD + i * (kw2 + Inches(0.24)), Inches(2.95), kw2, Inches(2.2), et, bas, gov, kenar=kenar)
not_serit(s, Inches(5.5),
          "Otomatik geçişler: son onarım kapandığında cihaz kendiliğinden 109 → 138 Ara Test'e geçer. "
          "İçine parça girilmemiş DGD işçilik satırı da diğer işler bitince otomatik tamamlanır.")

# ── 11 Test zinciri ────────────────────────────────────────────────────────
s = slayt(prs); eyebrow(s, "10", "Test zinciri")
baslik(s, "PhoneCheck entegrasyonu ve üç test noktası", pt=31)
lede(s, "Test verisi elle girilmiyor — IMEI okutulduğunda PhoneCheck API'sinden çekilip "
        "cihazın geçmişine yazılıyor.", Inches(1.95), w=Inches(6.9))
madde(s, PAD, Inches(2.95), Inches(6.7), Inches(1.9), [
    "Giriş Testi  103→104  — cihazı üretime taşıyan test",
    "Ara Test  138→124 / 138→109",
    "Son Test  125→126 / 125→109",
    "Onarım Bitiş Testi — Kamera / L3 / Ekran / Kasa için ayrı adım",
], pt=12)
yaz(s, PAD, Inches(4.95), Inches(6.7), Inches(0.9),
    [{"t": "Her test adımında en fazla 10 başarısız deneme hakkı var; hak dolunca sistem yeni "
           "denemeyi reddeder. Başarısızlar \"1. Başarısız … (1/10)\" diye numaralanır.",
      "pt": 11, "c": MUTED, "ara": 1.3}])
bx = Inches(7.85)
kutu(s, bx, Inches(2.4), Inches(4.75), Inches(3.3), dolgu=PANEL, kenar=LINE, radius=0.04)
yaz(s, bx + Inches(0.26), Inches(2.65), Inches(4.25), Inches(2.9),
    [{"t": "CİHAZ PHONECHECK'TE BULUNAMAZSA", "pt": 9.5, "c": SIGNAL, "kalin": True, "font": MONO, "alt": 7},
     {"t": "Test verisi elle doldurulur ama gerekçe ZORUNLUDUR — kayıt \"elle girildi\" olarak "
           "işaretlenir, denetlenebilir kalır.", "pt": 11, "c": MUTED, "ara": 1.3, "alt": 14},
     {"t": "VERİ EKSİK GELİRSE", "pt": 9.5, "c": SIGNAL, "kalin": True, "font": MONO, "alt": 7},
     {"t": "PhoneCheck cihazı bulup bazı alanları boş dönerse, gelenler dolu gelir, yalnız "
           "eksikler sorulur. Var olan kayıt tamamlanır — yeni deneme açılmaz, 10 hak boşuna "
           "tükenmez.", "pt": 11, "c": MUTED, "ara": 1.3}])

# ── 12 Fiyat ve onay ───────────────────────────────────────────────────────
s = slayt(prs); eyebrow(s, "11", "Fiyat ve onay")
baslik(s, "Onay kararını sistem veriyor")
lede(s, "Demontaj teknisyeni onarım planını kurar ve Üretime Aktar der. Sistem o anda parça "
        "fiyatları toplamını müşterinin hedef fiyat limiti ile karşılaştırır.", Inches(1.95),
     w=Inches(11))
kw3 = (ICER - Inches(0.26) * 2) / 3
kart(s, PAD, Inches(2.95), kw3, Inches(2.05), "LİMİT İÇİNDE", "Doğrudan üretime",
     "Cihaz 105 → 109 geçer, iş başlar. Müşteriye sorulmaz.", sol_serit=DONE)
kart(s, PAD + kw3 + Inches(0.26), Inches(2.95), kw3, Inches(2.05), "LİMİT AŞILDI", "Müşteri onayına",
     "Cihaz 105 → 106 gider. Buton \"Müşteri Onayı Al\"a döner ve basıldığında barkodlar da basılır.",
     sol_serit=SIGNAL)
kart(s, PAD + 2 * (kw3 + Inches(0.26)), Inches(2.95), kw3, Inches(2.05), "SONRADAN EKLENEN İŞ",
     "Karar yeniden verilir",
     "Üretimdeki cihaza yeni onarım eklenirse cihaz karar aşamasına geri çekilir ve toplam "
     "yeniden değerlendirilir.", sol_serit=LIVE)
yaz(s, PAD, Inches(5.35), ICER, Inches(0.9),
    [{"t": "Fiyat kaynakları: Müşteri Fiyat Matrisi (parça birim fiyatı) · Müşteri Hedef Fiyat "
           "Matrisi (cihaz başına üst limit) · Müşteri İşçilik Fiyatı Matrisi · seviye bazlı işçilik.",
      "pt": 11.5, "c": MUTED, "ara": 1.3}])

# ── 13 Depo ────────────────────────────────────────────────────────────────
s = slayt(prs); eyebrow(s, "12", "Depo")
baslik(s, "Parça, teslim edilmeden onarım kapanmıyor", pt=31)
madde(s, PAD, Inches(2.1), Inches(6.3), Inches(3.4), [
    "Depo → Parça Teslim ekranı: teknisyene çıkan her parça kaydedilir",
    "Sırası gelmemiş onarımın parçası teslim edilemez — ekranda \"Sırada Bekliyor, önce şu onarım\" görünür",
    "Reçete (Product BOM) cihazın modeline göre gösterilecek parçaları belirler",
    "Stok hareketleri, irsaliye ve iş emirleri aynı veriden beslenir",
    "DGD işçilik kalemi gerçek stok değildir — depo çıkışı yapılmaz",
], pt=12)
resim(s, "12_onarimhavuzu.png", Inches(7.15), Inches(2.15), Inches(5.45))
yaz(s, Inches(7.15), Inches(5.35), Inches(5.45), Inches(0.3),
    [{"t": "Onarım Havuzu — ekibin kendi iş listesi", "pt": 9.5, "c": MUTED, "font": MONO}])

# ── 14-15 Ekranlar ─────────────────────────────────────────────────────────
# İkişerli iki slayt: 2x2 yerleşimde görseller okunmayacak kadar küçülüyor ve
# alt sıra sayfadan taşıyordu. Görsel genişliği 5.9" -> yükseklik 3.23".
EKRAN_SLAYT = [
    ("13", "Ekranlar", "Sahada kullanılan arayüz", [
        ("07_technician.png", "Üretim Kaydını Görüntüle — onarım grupları, parçalar, statüler"),
        ("08_demontaj.png", "Üretime Aktar — onarım planı ve karar"),
    ]),
    ("13", "Ekranlar", "Sorgulama ve test ekranları", [
        ("03_servis.png", "Cihaz Sorgulama — statü geçmişi, testler, onarımlar"),
        ("10_sontest.png", "Son Test Sonuç — başarılıysa etiket basımı"),
    ]),
]
gw2 = Inches(5.9)
for kod, bolum, bas, gorseller in EKRAN_SLAYT:
    s = slayt(prs); eyebrow(s, kod, bolum)
    baslik(s, bas, pt=31)
    for i, (ad, alt) in enumerate(gorseller):
        x = PAD + i * (gw2 + Inches(0.3))
        resim(s, ad, x, Inches(2.35), gw2)
        yaz(s, x, Inches(5.78), gw2, Inches(0.34),
            [{"t": alt, "pt": 10, "c": MUTED, "font": MONO, "ara": 1.2}])

# ── 15 Etiket ve barkod ────────────────────────────────────────────────────
s = slayt(prs); eyebrow(s, "14", "Etiket ve barkod")
baslik(s, "Etiket tasarımı koda gömülü değil")
lede(s, "Şablonlar veritabanında tutulur ve Etiket Tasarımı ekranından düzenlenir. Üç tür "
        "etiket var: parça, üretime teslim, cihaz.", Inches(1.95), w=Inches(6.9))
madde(s, PAD, Inches(3.0), Inches(6.7), Inches(2.7), [
    "Barkod gerçek Code 39 yazı tipiyle basılır; yazı tipi hazır olmadan basım yapılmaz",
    "Başlangıç/bitiş işareti parantez — yıldız kaldırıldı",
    "DYMO 99014 (54×101 mm) için tam sayfa modu, kenar payı, \"tek sayfaya sığdır\" ölçümü",
    "Yazdır'a basınca pencere iş bitince kendi kapanır; aynı etiket iki kez basılamaz",
    "Demontaj'da \"Üretime Aktar\" ve \"Müşteri Onayı Al\" barkodları otomatik basar",
], pt=11.5)
ex = Inches(7.85)
kutu(s, ex, Inches(2.4), Inches(4.75), Inches(3.0), dolgu=PANEL, kenar=LINE, radius=0.04)
yaz(s, ex + Inches(0.26), Inches(2.62), Inches(4.25), Inches(0.3),
    [{"t": "ÜRETİME TESLİM ETİKETİ", "pt": 9.5, "c": SIGNAL, "kalin": True, "font": MONO}])
yaz(s, ex, Inches(3.2), Inches(4.75), Inches(0.8),
    [{"t": "( 3 5 8 1 2 0 )", "pt": 30, "c": TEXT, "kalin": True, "font": MONO}],
    hiza=PP_ALIGN.CENTER)
yaz(s, ex + Inches(0.26), Inches(4.15), Inches(4.25), Inches(1.0),
    [{"t": "IMEI · Model · Parti · Açıklama · Lokasyon · Referans", "pt": 10.5, "c": MUTED,
      "ara": 1.3, "alt": 8},
     {"t": "Form eksiksiz dolmadan yazdırılamaz.", "pt": 11, "c": SIGNAL, "kalin": True}],
    hiza=PP_ALIGN.CENTER)

# ── 16 Bu dönemde eklenenler ───────────────────────────────────────────────
s = slayt(prs); eyebrow(s, "15", "Bu dönemde eklenenler")
baslik(s, "Kurallar sisteme yerleşti", pt=31)
yeni = [
    ("SEVİYELENDİRME", "RMA → L3 → Batarya/Kamera/Kasa/Ekran → L1/L2 sırası kuruldu. Sırası gelmeyen iş 1004 ile bekletiliyor."),
    ("GRUP SEVİYESİNDE İŞLEM", "Tamamlama artık tek parça satırını değil onarımın tamamını kapatıyor."),
    ("OTOMATİK ARA TEST", "Son onarım kapanınca cihaz 109 → 138 kendiliğinden geçiyor; ayrı buton kaldırıldı."),
    ("PARÇA ÇAKIŞMA KURALLARI", "11 kural, 185 parça kategorisinin tamamında sınandı. Backend ve arayüz aynı sonucu veriyor."),
    ("TEST DÖNÜŞÜ", "Ara/son testten dönen cihazda yalnız L1/L2 teknisyene geri atanıyor."),
    ("TEK FORM", "Test verisi ve barkod bilgileri tek formda toplandı; eksik alan varken kaydedilemiyor."),
    ("AÇIK İŞ SAYACI", "Teknisyenin açık iş sayısı onarım bazında sayılıyor ve iş bitince düşüyor."),
    ("KAPANIŞ ZAMANI", "closed_at eklendi; statü geçmişinde işin gerçekten kapandığı an görünüyor."),
    ("CİHAZ İADE İZİ", "İade prosedürü statü geçmişine sebebiyle birlikte yazılıyor."),
]
kw4 = (ICER - Inches(0.24) * 2) / 3
for i, (et, gov) in enumerate(yeni):
    x = PAD + (i % 3) * (kw4 + Inches(0.24))
    y = Inches(2.0) + (i // 3) * Inches(1.62)
    kart(s, x, y, kw4, Inches(1.42), et, None, gov)

# ── 17 Sayılarla ───────────────────────────────────────────────────────────
s = slayt(prs); eyebrow(s, "16", "Sayılarla")
baslik(s, "Projenin bugünkü hâli")
sayilar1 = [("33.5K", "satır arayüz kodu\n(63 dosya)", TEXT), ("18.7K", "satır Python\niş mantığı", TEXT),
            ("48", "veri modeli\n(tablo eşleşmesi)", TEXT), ("185", "parça\nkategorisi", TEXT)]
sayilar2 = [("40+", "ekran / menü\ngirişi", SIGNAL), ("24", "cihaz\nstatüsü", SIGNAL),
            ("11", "parça çakışma\nkuralı", SIGNAL), ("10", "test denemesi\nüst sınırı", SIGNAL)]
kw5 = (ICER - Inches(0.24) * 3) / 4
for i, (v, k, c) in enumerate(sayilar1):
    istatistik(s, PAD + i * (kw5 + Inches(0.24)), Inches(2.2), kw5, Inches(1.62), v, k, c)
for i, (v, k, c) in enumerate(sayilar2):
    istatistik(s, PAD + i * (kw5 + Inches(0.24)), Inches(4.05), kw5, Inches(1.62), v, k, c)
yaz(s, PAD, Inches(6.0), ICER, Inches(0.6),
    [{"t": "Veritabanı: PostgreSQL — warehouse ve organization şemaları, Türkçe sıralama için "
           "ICU tr-TR collation.", "pt": 11.5, "c": MUTED, "ara": 1.3}])

# ── 18 Sıradaki adımlar ────────────────────────────────────────────────────
s = slayt(prs); eyebrow(s, "17", "Sıradaki adımlar")
baslik(s, "Açık başlıklar")
kw6 = (ICER - Inches(0.3)) / 2
kutu(s, PAD, Inches(2.2), kw6, Inches(2.6), dolgu=PANEL, kenar=LINE, radius=0.05)
yaz(s, PAD + Inches(0.26), Inches(2.44), kw6 - Inches(0.5), Inches(0.3),
    [{"t": "VERİ TEMİZLİĞİ", "pt": 9.5, "c": SIGNAL, "kalin": True, "font": MONO}])
madde(s, PAD + Inches(0.26), Inches(2.9), kw6 - Inches(0.5), Inches(1.7), [
    "Kurallar öncesinde girilmiş Back Cover + Back Glass çakışmaları (10 cihaz) ve fazla batarya flexleri (5 cihaz)",
    "Cihaza bağlanamayan 75 sahipsiz onarım grubu (209 satır)",
], pt=11.5)
x3 = PAD + kw6 + Inches(0.3)
kutu(s, x3, Inches(2.2), kw6, Inches(2.6), dolgu=PANEL, kenar=LINE, radius=0.05)
yaz(s, x3 + Inches(0.26), Inches(2.44), kw6 - Inches(0.5), Inches(0.3),
    [{"t": "EKSİK / KARAR BEKLEYEN", "pt": 9.5, "c": SIGNAL, "kalin": True, "font": MONO}])
madde(s, x3 + Inches(0.26), Inches(2.9), kw6 - Inches(0.5), Inches(1.7), [
    "Rezervasyon sistemi henüz yok",
    "Hedef fiyat matrisinin kapsamı genişletilmeli",
    "mission_groups.order_number için arayüzde düzenleme ekranı yok",
], pt=11.5)
not_serit(s, Inches(5.15),
          "Parça çakışma kuralları ve seviyelendirme artık sistemin çekirdeğinde. Bundan sonraki "
          "kurallar tek satır ekleyerek tanımlanabiliyor — liste backend ve arayüzde tek yerde.")

# ── 19 Kapanış ─────────────────────────────────────────────────────────────
s = slayt(prs)
kutu(s, 0, 0, Inches(0.09), H, dolgu=SIGNAL)
yaz(s, PAD, Inches(2.3), ICER, Inches(0.4),
    [{"t": "SON", "pt": 12, "c": SIGNAL, "kalin": True, "font": MONO}])
yaz(s, PAD, Inches(2.8), Inches(9.5), Inches(1.2),
    [{"t": "Sorular?", "pt": 52, "c": TEXT, "kalin": True}])
lede(s, "RemaLab WMS — depo, servis ve onarım operasyonunun tamamı tek sistemde.",
     Inches(4.25), pt=15, w=Inches(9.2))
kutu(s, PAD, Inches(5.1), Inches(3.4), Emu(28575), dolgu=SIGNAL)

prs.save(CIKTI)
print(f"kaydedildi: {CIKTI}")
print(f"slayt sayisi: {len(prs.slides.__iter__.__self__._sldIdLst)}")
print(f"boyut: {os.path.getsize(CIKTI) // 1024} KB")
