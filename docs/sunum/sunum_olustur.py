# -*- coding: utf-8 -*-
"""RemaLab WMS - PowerPoint sunumu üretici.

Tüm rakamlar canlı veritabanından ve kaynak koddan alınmıştır.
Ekran görüntüleri capture.py ile gerçek uygulamadan yakalanmıştır.
"""
import os
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

BASE = r"C:\Users\DELL\Desktop\remalab-feature-users-module\remalab-feature-users-module"
SHOTS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ekran_goruntuleri")
OUT = os.path.join(BASE, "docs", "RemaLab_WMS_Sunum.pptx")

# ── Palet (uygulamanın kendi koyu teması) ────────────────────────────
BG        = RGBColor(0x09, 0x0A, 0x0F)
BG2       = RGBColor(0x0E, 0x10, 0x17)
PANEL     = RGBColor(0x12, 0x14, 0x1C)
PANEL2    = RGBColor(0x1A, 0x1D, 0x28)
BORDER    = RGBColor(0x1E, 0x22, 0x2D)
BORDER2   = RGBColor(0x2E, 0x35, 0x45)
CYAN      = RGBColor(0x00, 0xB2, 0xFF)
VIOLET    = RGBColor(0x7C, 0x5C, 0xFF)
GREEN     = RGBColor(0x00, 0xE6, 0x76)
AMBER     = RGBColor(0xFF, 0x95, 0x00)
RED       = RGBColor(0xFF, 0x3B, 0x30)
ROSE      = RGBColor(0xC2, 0x44, 0x5F)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT      = RGBColor(0xE8, 0xEC, 0xF5)
MUTED     = RGBColor(0x88, 0x92, 0xB5)
DIM       = RGBColor(0x5A, 0x66, 0x85)

FONT = "Segoe UI"
FONT_MONO = "Consolas"

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

_sayac = {"n": 0}


# ── Yardımcılar ──────────────────────────────────────────────────────
def _alpha(shape, pct):
    """Dolgu saydamlığı. python-pptx bunu API'den vermiyor, XML'e eklenir."""
    sp = shape.fill._xPr
    srgb = sp.find(qn('a:solidFill'))
    if srgb is None:
        return
    clr = srgb.find(qn('a:srgbClr'))
    if clr is None:
        return
    a = etree.SubElement(clr, qn('a:alpha'))
    a.set('val', str(int(pct * 1000)))


def rect(slide, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE,
         alpha=None, adj=None):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        if alpha is not None:
            _alpha(s, alpha)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    if adj is not None:
        try:
            s.adjustments[0] = adj
        except Exception:
            pass
    s.text_frame.text = ""
    return s


def grad(slide, x, y, w, h, c1, c2, angle=0, alpha1=None, alpha2=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    s.line.fill.background()
    s.fill.gradient()
    st = s.fill.gradient_stops
    st[0].color.rgb = c1
    st[0].position = 0.0
    st[1].color.rgb = c2
    st[1].position = 1.0
    s.fill.gradient_angle = angle
    if alpha1 is not None or alpha2 is not None:
        gs = s.fill._xPr.find(qn('a:gradFill')).find(qn('a:gsLst'))
        for node, a in zip(gs.findall(qn('a:gs')), [alpha1, alpha2]):
            if a is None:
                continue
            clr = node.find(qn('a:srgbClr'))
            el = etree.SubElement(clr, qn('a:alpha'))
            el.set('val', str(int(a * 1000)))
    return s


def txt(slide, x, y, w, h, text, size=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
        font=FONT, spacing=None, anchor=MSO_ANCHOR.TOP, italic=False, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        if spacing is not None:
            r.font._rPr.set('spc', str(int(spacing * 100)))
    return tb


def bullets(slide, x, y, w, h, items, size=13, color=TEXT, gap=8, bullet_color=CYAN):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.25
        r1 = p.add_run()
        r1.text = "▸  "
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r1.font.bold = True
        r1.font.name = FONT
        r2 = p.add_run()
        r2.text = it
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = FONT
    return tb


def bg(slide, tone=BG):
    rect(slide, 0, 0, 13.333, 7.5, fill=tone)


def glow(slide, cx, cy, r, color, alpha=8):
    """Yumuşak ışık hissi: iç içe geçen saydam daireler."""
    for i, k in enumerate([1.0, 0.66, 0.4]):
        rr = r * k
        rect(slide, cx - rr, cy - rr, rr * 2, rr * 2, fill=color,
             shape=MSO_SHAPE.OVAL, alpha=alpha * (i + 1) * 0.6)


def footer(slide, label=""):
    _sayac["n"] += 1
    rect(slide, 0, 7.38, 13.333, 0.12, fill=PANEL)
    rect(slide, 0, 7.38, 13.333 * 0.28, 0.12, fill=CYAN)
    txt(slide, 0.6, 7.02, 6, 0.3, "REMALAB WMS", 8, DIM, True, spacing=2.2)
    if label:
        txt(slide, 3.0, 7.02, 7, 0.3, label.upper(), 8, DIM, False, spacing=1.6)
    txt(slide, 11.5, 7.02, 1.25, 0.3, "%02d" % _sayac["n"], 9, MUTED, True,
        align=PP_ALIGN.RIGHT, font=FONT_MONO)


def new(tone=BG):
    s = prs.slides.add_slide(BLANK)
    bg(s, tone)
    return s


def chip(slide, x, y, text, color=CYAN, w=None, size=9):
    ww = w or (0.16 * len(text) + 0.34)
    rect(slide, x, y, ww, 0.32, fill=color, alpha=16,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.5, line=color, lw=0.75)
    txt(slide, x, y + 0.055, ww, 0.24, text, size, color, True, align=PP_ALIGN.CENTER, spacing=1.2)
    return ww


def title_block(slide, kicker, title, sub=None, kcolor=CYAN):
    chip(slide, 0.75, 0.62, kicker, kcolor)
    txt(slide, 0.75, 1.08, 11.8, 0.7, title, 32, WHITE, True)
    if sub:
        txt(slide, 0.75, 1.83, 11.5, 0.5, sub, 13, MUTED)
    rect(slide, 0.75, 2.32, 1.5, 0.045, fill=kcolor)


def card(slide, x, y, w, h, fill=PANEL, line=BORDER, accent=None, aw=0.05):
    c = rect(slide, x, y, w, h, fill=fill, line=line, lw=0.9,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.045)
    if accent:
        rect(slide, x, y, aw, h, fill=accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.5)
    return c


def kpi(slide, x, y, w, h, value, label, note=None, color=CYAN):
    card(slide, x, y, w, h, accent=color)
    txt(slide, x + 0.32, y + 0.34, w - 0.5, 0.7, value, 30, WHITE, True)
    txt(slide, x + 0.32, y + 1.02, w - 0.5, 0.3, label, 10.5, MUTED, True, spacing=0.8)
    if note:
        txt(slide, x + 0.32, y + 1.32, w - 0.5, 0.3, note, 9, DIM)


def pic(slide, path, x, y, w, border=True, shadow_pad=0.0):
    """Görsel + ince çerçeve. Yükseklik en-boy oranından hesaplanır."""
    from PIL import Image
    iw, ih = Image.open(path).size
    h = w * ih / iw
    if border:
        rect(slide, x - 0.035, y - 0.035, w + 0.07, h + 0.07, fill=BORDER2,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.02)
    slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    return h


def shot(path):
    return os.path.join(SHOTS, path)


# ════════════════════════════════════════════════════════════════════
# 1 — KAPAK
# ════════════════════════════════════════════════════════════════════
s = new(BG)
grad(s, 0, 0, 13.333, 7.5, RGBColor(0x0B, 0x10, 0x22), BG, angle=45)
glow(s, 11.2, 1.4, 3.4, CYAN, 7)
glow(s, 1.4, 6.6, 2.8, VIOLET, 6)
# ince ızgara hissi
for i in range(1, 14):
    rect(s, i * 0.95, 0, 0.008, 7.5, fill=WHITE, alpha=2.5)

logo = os.path.join(BASE, "frontend", "public", "karanlık-mod.png")
if os.path.exists(logo):
    from PIL import Image as _I
    lw, lh = _I.open(logo).size
    s.shapes.add_picture(logo, Inches(0.95), Inches(0.75), Inches(2.9), Inches(2.9 * lh / lw))

chip(s, 0.98, 3.05, "WAREHOUSE MANAGEMENT SYSTEM", CYAN)
txt(s, 0.95, 3.55, 11.5, 1.0, "RemaLab WMS", 60, WHITE, True)
txt(s, 0.95, 4.62, 10.5, 0.9,
    "Cihaz yaşam döngüsünü kabulden sevkiyata kadar tek sistemde yöneten\n"
    "masaüstü depo ve onarım yönetim platformu",
    16, MUTED, line_spacing=1.35)

rect(s, 0.98, 5.75, 1.9, 0.05, fill=CYAN)

for i, (k, v) in enumerate([("30.126", "parça çeşidi"), ("7.679", "cihaz kaydı"),
                            ("194", "backend servisi"), ("35", "ekran")]):
    xx = 0.95 + i * 2.35
    txt(s, xx, 6.05, 2.2, 0.45, k, 22, CYAN, True)
    txt(s, xx, 6.52, 2.2, 0.3, v, 10, DIM, True, spacing=1.2)

txt(s, 9.6, 6.85, 3.0, 0.3, "Ağustos 2026  ·  v1.0.0", 10, DIM, align=PP_ALIGN.RIGHT)
_sayac["n"] = 1

# ════════════════════════════════════════════════════════════════════
# 2 — AJANDA
# ════════════════════════════════════════════════════════════════════
s = new()
glow(s, 12.4, 0.6, 2.6, CYAN, 5)
title_block(s, "İÇİNDEKİLER", "Sunum Akışı",
            "Sistemin ne olduğundan mühendislik detaylarına kadar altı bölüm")

items = [
    ("01", "Sistem Nedir?", "Amaç, mimari, teknoloji yığını", CYAN),
    ("02", "Rakamlarla RemaLab", "Canlı veritabanından ölçümler", GREEN),
    ("03", "Cihaz Yaşam Döngüsü", "23 statü, 27 geçiş, 17 görev grubu", AMBER),
    ("04", "Ekranlar", "18 ekran görüntüsüyle uçtan uca tur", VIOLET),
    ("05", "Mühendislik", "QWebChannel köprüsü, performans, tema", ROSE),
    ("06", "Kalite ve Yol Haritası", "Düzeltilen hatalar, sıradaki adımlar", RED),
]
for i, (no, t, d, c) in enumerate(items):
    col, row = i % 2, i // 2
    x = 0.75 + col * 6.2
    y = 2.72 + row * 1.42
    card(s, x, y, 5.85, 1.2, accent=c)
    txt(s, x + 0.32, y + 0.26, 0.8, 0.5, no, 22, c, True, font=FONT_MONO)
    txt(s, x + 1.25, y + 0.24, 4.3, 0.35, t, 15, WHITE, True)
    txt(s, x + 1.25, y + 0.66, 4.3, 0.35, d, 10.5, MUTED)
footer(s, "İçindekiler")


# ── Bölüm ayracı yardımcısı ─────────────────────────────────────────
def section(no, title, sub, color=CYAN):
    sl = new(BG2)
    grad(sl, 0, 0, 13.333, 7.5, BG, RGBColor(0x0D, 0x11, 0x1D), angle=30)
    glow(sl, 10.8, 5.6, 3.6, color, 7)
    glow(sl, 2.2, 1.4, 2.2, color, 4)
    txt(sl, 0.9, 1.55, 4, 2.6, no, 150, color, True, font=FONT_MONO)
    _alpha_box = rect(sl, 0.95, 4.02, 2.4, 0.06, fill=color)
    txt(sl, 0.95, 4.35, 11, 0.9, title, 40, WHITE, True)
    txt(sl, 0.95, 5.35, 9.5, 0.6, sub, 14, MUTED)
    footer(sl, title)
    return sl


# ════════════════════════════════════════════════════════════════════
section("01", "Sistem Nedir?", "Amaç, mimari ve teknoloji yığını")

# 3 — NE İŞE YARAR
s = new()
title_block(s, "GENEL BAKIŞ", "RemaLab WMS ne yapar?",
            "Yenilenmiş cihaz operasyonunun tamamını tek bir masaüstü uygulamada toplar")
kartlar = [
    ("▦", "Depo & Stok", "30 binden fazla parça çeşidi, çoklu lokasyon, kritik stok uyarıları ve "
                         "Good Stock / Repair Stock ayrımı.", CYAN),
    ("↻", "Cihaz Yaşam Döngüsü", "Kayıt kabulden müşteri sevkiyatına kadar 23 statü ve "
                                 "27 tanımlı geçiş; her adım denetlenebilir.", GREEN),
    ("⚒", "Onarım Yönetimi", "Departman bazlı onarım havuzları, teknisyen atama, parça talebi ve "
                             "müşteri onayı akışı.", AMBER),
    ("✓", "Test Entegrasyonu", "PhoneCheck test sonuçları otomatik çekilir; kritik parça "
                               "orijinalliği (Genuine / Not Genuine) izlenir.", VIOLET),
]
for i, (ikon, t, d, c) in enumerate(kartlar):
    col, row = i % 2, i // 2
    x = 0.75 + col * 6.2
    y = 2.72 + row * 2.1
    card(s, x, y, 5.85, 1.9, accent=c)
    rect(s, x + 0.34, y + 0.36, 0.42, 0.42, fill=c, alpha=18,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.25, line=c, lw=0.75)
    txt(s, x + 0.34, y + 0.44, 0.42, 0.3, ikon, 13, c, True, align=PP_ALIGN.CENTER)
    txt(s, x + 1.0, y + 0.34, 4.5, 0.4, t, 15, WHITE, True)
    txt(s, x + 1.0, y + 0.82, 4.55, 1.0, d, 10.5, MUTED, line_spacing=1.3)
footer(s, "Genel Bakış")

# 4 — MİMARİ
s = new()
title_block(s, "MİMARİ", "Tek süreç, üç katman",
            "Ağ katmanı yok — React ile Python aynı süreçte QWebChannel üzerinden konuşur")

# Katman kutuları — (y, yükseklik, başlık, teknoloji, not, renk)
kats = [
    (2.52, 1.02, "SUNUM KATMANI", "React 19  ·  Vite 8  ·  Tailwind 4", "35 ekran · 27.354 satır", CYAN),
    (3.92, 0.92, "KÖPRÜ", "QWebChannel  (194 @Slot)", "JSON string · çift yönlü", VIOLET),
    (5.22, 1.02, "İŞ MANTIĞI", "PySide6  ·  SQLAlchemy  ·  14 servis", "15.781 satır Python", GREEN),
]
for y, hh, t, d, n, c in kats:
    card(s, 0.9, y, 8.4, hh, accent=c, aw=0.06)
    txt(s, 1.35, y + 0.14, 6, 0.28, t, 10, c, True, spacing=1.8)
    txt(s, 1.35, y + 0.44, 7, 0.34, d, 14, WHITE, True)
    txt(s, 1.35, y + 0.76, 7, 0.26, n, 9.5, DIM)

# oklar
for yy in (3.58, 4.88):
    txt(s, 4.85, yy, 0.5, 0.3, "▼", 14, BORDER2, True, align=PP_ALIGN.CENTER)
txt(s, 4.85, 6.28, 0.5, 0.3, "▼", 14, BORDER2, True, align=PP_ALIGN.CENTER)

# Veritabanı
card(s, 0.9, 6.52, 8.4, 0.5, fill=PANEL2, line=BORDER)
txt(s, 1.35, 6.64, 7.5, 0.28, "PostgreSQL   ·   70 tablo   ·   warehouse + organization şemaları",
    11.5, TEXT, True)

# Sağ bilgi paneli
card(s, 9.75, 2.52, 2.85, 4.5, fill=PANEL2, line=BORDER)
txt(s, 10.05, 2.78, 2.3, 0.3, "NEDEN BÖYLE?", 10, CYAN, True, spacing=1.6)
bullets(s, 10.05, 3.2, 2.3, 3.6, [
    "REST API yok — HTTP gecikmesi ve serileştirme maliyeti ortadan kalkıyor",
    "Tek dağıtım paketi; kullanıcı tek exe/klasör çalıştırıyor",
    "Web arayüzünün esnekliği + masaüstü uygulamanın donanım erişimi",
    "Barkod okuyucu doğrudan input'a yazıyor",
], size=9.5, gap=9)
footer(s, "Mimari")

# 5 — TEKNOLOJİ
s = new()
title_block(s, "TEKNOLOJİ", "Yığın", "Üretimde kullanılan sürümler")

sol = [("React", "19.2"), ("React Router", "7.18"), ("Vite", "8.1"),
       ("Tailwind CSS", "4.3"), ("lucide-react", "1.24"), ("SheetJS (xlsx)", "0.18")]
sag = [("Python", "3.13"), ("PySide6 / Qt", "6.11"), ("SQLAlchemy", "2.x"),
       ("PostgreSQL", "sunucu"), ("psycopg2", "2.x"), ("PhoneCheck API", "entegrasyon")]

for i, (baslik, liste, renk) in enumerate([("FRONTEND", sol, CYAN), ("BACKEND", sag, GREEN)]):
    x = 0.75 + i * 6.2
    card(s, x, 2.72, 5.85, 4.0, accent=renk)
    txt(s, x + 0.4, 2.98, 4, 0.32, baslik, 11, renk, True, spacing=2.0)
    for j, (ad, ver) in enumerate(liste):
        yy = 3.5 + j * 0.5
        rect(s, x + 0.4, yy + 0.13, 0.07, 0.07, fill=renk, shape=MSO_SHAPE.OVAL)
        txt(s, x + 0.68, yy, 3.2, 0.32, ad, 12.5, TEXT)
        txt(s, x + 3.9, yy, 1.5, 0.32, ver, 11, DIM, align=PP_ALIGN.RIGHT, font=FONT_MONO)

txt(s, 0.75, 6.92, 12, 0.3,
    "Kullanılmayan bağımlılıklar tespit edildi: axios, zustand, qwebchannel (npm) — temizlenebilir.",
    9.5, DIM, italic=True)
footer(s, "Teknoloji Yığını")

# ════════════════════════════════════════════════════════════════════
section("02", "Rakamlarla RemaLab", "Canlı veritabanından alınan ölçümler", GREEN)

# 6 — KPI
s = new()
title_block(s, "CANLI VERİ", "Sistem bugün ne kadar veri taşıyor?",
            "Rakamlar üretim veritabanından okundu — 4 Ağustos 2026", GREEN)
kpis = [
    ("30.126", "PARÇA ÇEŞİDİ", "30.187 stok satırı", CYAN),
    ("7.679", "CİHAZ KAYDI", "batch_entries", GREEN),
    ("159", "KULLANICI", "8 farklı rol", AMBER),
    ("70", "VERİTABANI TABLOSU", "2 şema", VIOLET),
    ("23", "STATÜ TANIMI", "27 aktif geçiş", ROSE),
    ("17", "GÖREV GRUBU", "33 görev tanımı", RED),
]
for i, (v, l, n, c) in enumerate(kpis):
    col, row = i % 3, i // 3
    kpi(s, 0.75 + col * 4.13, 2.72 + row * 2.05, 3.85, 1.8, v, l, n, c)
footer(s, "Rakamlarla")

# 7 — KOD TABANI
s = new()
title_block(s, "KOD TABANI", "Proje büyüklüğü", "Satır sayıları ve dağılım", GREEN)

cols = [("Frontend", "27.354", "satır JSX/JS", CYAN,
         ["35 sayfa bileşeni", "12 ortak bileşen", "191 API sarmalayıcı", "389 satır tema CSS"]),
        ("Backend", "15.781", "satır Python", GREEN,
         ["194 @Slot metodu", "44 SQLAlchemy modeli", "14 servis sınıfı", "8 repository"]),
        ("En büyük dosyalar", "3.426", "satır (WorkOrders.jsx)", AMBER,
         ["web_bridge.py — 12.121", "BatchEntry.jsx — 1.817", "api.js — 2.430",
          "TechnicianRepair — 1.479"])]

for i, (t, v, u, c, lst) in enumerate(cols):
    x = 0.75 + i * 4.13
    card(s, x, 2.72, 3.85, 4.0, accent=c)
    txt(s, x + 0.35, 3.0, 3.2, 0.3, t.upper(), 10, c, True, spacing=1.6)
    txt(s, x + 0.35, 3.38, 3.2, 0.6, v, 28, WHITE, True)
    txt(s, x + 0.35, 4.0, 3.2, 0.3, u, 10, DIM)
    rect(s, x + 0.35, 4.42, 3.15, 0.03, fill=BORDER2)
    bullets(s, x + 0.35, 4.65, 3.2, 2.0, lst, size=10, gap=7, bullet_color=c)
footer(s, "Kod Tabanı")

# ════════════════════════════════════════════════════════════════════
section("03", "Cihaz Yaşam Döngüsü", "Kayıt kabulden sevkiyata kadar izlenen yol", AMBER)

# 8 — STATÜ AKIŞI
s = new()
title_block(s, "AKIŞ", "Cihaz hangi aşamalardan geçiyor?",
            "Her geçiş service_statu_map tablosunda tanımlı ve denetleniyor", AMBER)

akis = [
    ("100", "Kayıt\nKabul", CYAN), ("101", "Depo\nGirişi", CYAN),
    ("102", "İlk Test\nTransferi", VIOLET), ("103", "İlk Test\nTamam", VIOLET),
    ("104", "Üretime\nTeslim", GREEN), ("105", "Teknik\nKabul", GREEN),
    ("109", "Üretim\nDevam", AMBER), ("138", "Ara Test\nBekliyor", ROSE),
    ("124", "Son Test\nTransferi", ROSE), ("125", "Son Test\nKabul", ROSE),
    ("126", "Test\nBaşarılı", GREEN), ("127", "Sevkiyat\nBekliyor", CYAN),
]
CW, GAP = 1.66, 0.37       # kart genisligi + aradaki ok bosluğu
for i, (kod, ad, c) in enumerate(akis):
    col, row = i % 6, i // 6
    x = 0.75 + col * (CW + GAP)
    y = 2.62 + row * 1.62
    card(s, x, y, CW, 1.24, fill=PANEL, accent=c, aw=0.05)
    txt(s, x + 0.22, y + 0.18, 1.3, 0.35, kod, 17, c, True, font=FONT_MONO)
    txt(s, x + 0.22, y + 0.58, 1.3, 0.58, ad, 9.5, MUTED, line_spacing=1.15)
    if col < 5:
        txt(s, x + CW, y + 0.44, GAP, 0.34, "›", 22, BORDER2, True, align=PP_ALIGN.CENTER)
    elif row == 0:
        # satir sonu: akis alttaki satirin basina donuyor
        txt(s, x + CW - 0.1, y + 1.28, GAP, 0.3, "↴", 16, BORDER2, True, align=PP_ALIGN.CENTER)

card(s, 0.75, 5.92, 11.85, 1.05, fill=PANEL2, line=BORDER, accent=AMBER)
txt(s, 1.12, 6.12, 5, 0.3, "CANLI DAĞILIM", 9.5, AMBER, True, spacing=1.6)
txt(s, 1.12, 6.46, 11, 0.3,
    "102 → 2.258 cihaz     ·     109 → 2.221     ·     127 → 1.137     ·     "
    "104 → 828     ·     101 → 602     ·     125 → 287",
    11, MUTED)
footer(s, "Cihaz Yaşam Döngüsü")

# 9 — DEPARTMANLAR
s = new()
title_block(s, "ORGANİZASYON", "Departmanlar ve onarım havuzları",
            "17 görev grubu, 33 görev tanımı — yetki bu eşleşmeden okunuyor", AMBER)

deps = [("BATTERY", "Batarya Onarımı", GREEN), ("CAMERA", "Kamera Onarımı", CYAN),
        ("DISPLAY", "Ekran Onarımı", VIOLET), ("CASE", "Kasa Onarımı", AMBER),
        ("L1REPAIR", "L1 Onarımı", ROSE), ("L2REPAIR", "L2 Onarımı", ROSE),
        ("L3REPAIR", "L3 Onarımı", RED), ("DISMANTLE", "Demontaj", MUTED)]
for i, (kod, ad, c) in enumerate(deps):
    col, row = i % 4, i // 4
    x = 0.75 + col * 3.1
    y = 2.75 + row * 1.35
    card(s, x, y, 2.85, 1.15, accent=c)
    txt(s, x + 0.28, y + 0.22, 2.3, 0.3, kod, 10, c, True, font=FONT_MONO, spacing=1.2)
    txt(s, x + 0.28, y + 0.58, 2.4, 0.35, ad, 13, WHITE, True)

card(s, 0.75, 5.62, 11.85, 1.2, fill=PANEL2, line=BORDER)
txt(s, 1.1, 5.85, 11, 0.3, "ONARIM DURUM KODLARI", 10, AMBER, True, spacing=1.6)
kodlar = "1000 Teknisyene Atanacak   ·   1001 Teknisyene Atandı   ·   1002 Tamamlandı   ·   " \
         "1003 İptal   ·   1004 Yüksek Seviye   ·   1005 Müşteri Onayı   ·   1006 Test Bekleniyor"
txt(s, 1.1, 6.22, 11, 0.4, kodlar, 10.5, MUTED)
footer(s, "Departmanlar")

# ════════════════════════════════════════════════════════════════════
section("04", "Ekranlar", "Uygulamanın 18 ekran görüntüsüyle uçtan uca turu", VIOLET)


def ekran(baslik, aciklama, dosya, notlar, renk=VIOLET, etiket="EKRAN"):
    sl = new()
    glow(sl, 12.6, 0.4, 2.4, renk, 4)
    chip(sl, 0.75, 0.55, etiket, renk)
    txt(sl, 0.75, 0.98, 8.2, 0.5, baslik, 25, WHITE, True)
    txt(sl, 0.75, 1.52, 8.2, 0.4, aciklama, 11, MUTED)
    pic(sl, shot(dosya), 0.75, 2.05, 8.35)
    card(sl, 9.45, 2.05, 3.15, 4.6, fill=PANEL2, line=BORDER, accent=renk)
    txt(sl, 9.78, 2.32, 2.6, 0.3, "ÖNE ÇIKANLAR", 9.5, renk, True, spacing=1.6)
    bullets(sl, 9.78, 2.75, 2.55, 3.7, notlar, size=9.5, gap=9, bullet_color=renk)
    footer(sl, baslik)
    return sl


ekran("Giriş Ekranı",
      "Kurumsal kimlikle karşılayan, animasyonlu mozaik zeminli giriş sayfası",
      "01_login.png",
      ["Kullanıcı adı + şifre doğrulaması backend'de bcrypt ile",
       "\"Beni Hatırla\" seçilirse oturum localStorage'da kalıcı",
       "Role göre otomatik yönlendirme: depo → /depo, diğerleri → /dashboard",
       "Uygulamanın tek animasyonlu dokulu ekranı"],
      CYAN, "01 · GİRİŞ")

ekran("Kontrol Paneli",
      "Depo operasyonunun canlı özeti — stok, kritik uyarılar ve son hareketler",
      "02_dashboard.png",
      ["5 canlı KPI kartı: parça çeşidi, kritik stok, giriş/çıkış, lokasyon",
       "Son stok hareketleri tablosu (Good Stock ⇄ Repair Stock transferleri)",
       "Hızlı erişim kısayolları",
       "Bildirim zili 60 saniyede bir kritik stoğu sorguluyor"],
      CYAN, "02 · KONTROL PANELİ")

ekran("Servis — Cihaz Sorgulama",
      "IMEI ile cihazın tüm geçmişi: PhoneCheck testleri ve kritik parça orijinalliği",
      "03_servis.png",
      ["9 sütunlu PhoneCheck test geçmişi",
       "Kritik parça orijinalliği: 🟩 Genuine / 🟥 Not Genuine / ⬜ Belirtilmemiş",
       "Ana Kamera, Batarya ve Eski Pil ayrı ayrı izleniyor",
       "Grade, BatteryCycle, StationID ve hata detayları"],
      ROSE, "03 · SERVİS")

ekran("Depo Yönetimi",
      "Lokasyon bazlı stok görünümü ve transfer işlemleri",
      "04_depo.png",
      ["Good Stock / Repair Stock ayrımı",
       "6 aktif lokasyon arasında transfer",
       "Kritik seviye altındaki parçalar renkli işaretleniyor",
       "Excel'e aktarma desteği"],
      AMBER, "04 · DEPO")

ekran("Parça Kartları",
      "30.126 parça çeşidinin merkezi kataloğu",
      "05_parts.png",
      ["Kategori, ürün ailesi ve tedarikçi bazlı filtreleme",
       "Stok takibi olan / olmayan parça ayrımı",
       "Excel ile toplu içe aktarma (sütun eşleme modali)",
       "Kritik stok eşiği parça bazında tanımlanıyor"],
      VIOLET, "05 · PARÇALAR")

ekran("İş Emirleri",
      "8 sekmeli üretim merkezi — projenin en kapsamlı ekranı (3.426 satır)",
      "06_workorders.png",
      ["Yarı Mamul Üretimi (varsayılan sekme)",
       "Servis iş emirleri, malzeme tüketimi, üretim raporu",
       "Barkod parça takip modülü gömülü",
       "Reçete (BOM) bazlı otomatik hammadde düşümü"],
      CYAN, "06 · İŞ EMİRLERİ")

ekran("Üretim Kaydını Görüntüle",
      "Teknisyenin cihaz üzerindeki tüm onarım kalemlerini yönettiği ekran",
      "07_technician.png",
      ["Onarım ekleme, teknisyen atama ve parça talebi",
       "Depo durumu 11 farklı tedarik statüsünden okunuyor",
       "Kritik parça orijinallik paneli",
       "Tamamlama üç kuralla korunuyor: atama + parça çıkışı + müşteri onayı"],
      RED, "07 · ÜRETİM TEKNİSYENİ")

ekran("Üretime Aktar (Demontaj)",
      "Demontaj teknisyeninin cihazı üretime devrettiği karar ekranı",
      "08_demontaj.png",
      ["Seçili onarıma parça ekleme butonu",
       "Atanmış teknisyen sütunu",
       "Karar butonu eklenen parçalara göre 109 veya 106 statüsünü seçiyor",
       "Görev grubu kilitli — yanlış departmana kayıt açılamıyor"],
      GREEN, "08 · DEMONTAJ")

ekran("Onarım Havuzu",
      "Departman bazlı onarım kuyruğu — tek bileşen, 7 departman",
      "12_onarimhavuzu.png",
      ["Rota parametresi departmanı belirliyor (/onarim-havuzu/BATTERY)",
       "Kod kopyalanmadığı için düzeltmeler tüm departmanlara aynı anda iniyor",
       "Teknisyen kendi görev grubundaki kayıtları görüyor",
       "BATTERY, CAMERA, DISPLAY, CASE, L1/L2/L3"],
      AMBER, "09 · ONARIM HAVUZU")

ekran("Hızlı Onarım Bitiş",
      "Tek barkod kutusu, buton yok — okut ve kapat",
      "09_hizli_batarya.png",
      ["Barkod okuyucu doğrudan input'a yazıyor, Enter ile işleniyor",
       "Her kayıt için renk kodlu sonuç satırı",
       "Kısmi kapanma destekli: 4 kayıttan 3'ü kapanabilir",
       "Tamamlama kuralları \"Onarımı Tamamla\" ile ortak kod paylaşıyor"],
      RED, "10 · HIZLI BİTİŞ")

ekran("Kayıt Kabul (100 → 101)",
      "Statü geçiş ekranı — tek bileşen 9 farklı menü öğesine hizmet ediyor",
      "11_kayitkabul.png",
      ["Geçiş tanımı veritabanından okunuyor, kodda sabit değil",
       "Kaynak statü tutmuyorsa işlem reddediliyor",
       "IMEI / seri no / internal ID / batch no ile arama",
       "Renk kodlu işlem geçmişi log'u"],
      CYAN, "11 · STATÜ GEÇİŞİ")

ekran("Son Test Sonuç",
      "Başarılı ise depoya, başarısız ise hata kodlarıyla tekniğe geri",
      "10_sontest.png",
      ["Sol sütun onay, sağ sütun geri çevrim",
       "Hata kataloğundan en az 1, en fazla 10 arıza seçimi",
       "Açıklama zorunlu",
       "Ara Test ekranı aynı bileşeni farklı statülerle kullanıyor"],
      VIOLET, "12 · SON TEST")

ekran("Ara Test Yap (138 → 124)",
      "Üretim sonrası ara kalite kontrolü",
      "19_aratest.png",
      ["Başarılı → 124 (son teste transfer)",
       "Başarısız → 109 (üretime geri)",
       "TestResultScreen bileşeninin ikinci kullanımı",
       "13 satırlık sarmalayıcı ile tanımlanıyor"],
      ROSE, "13 · ARA TEST")

ekran("Kullanıcı Yönetimi",
      "159 kullanıcı, 8 rol ve görev grubu ataması",
      "13_users.png",
      ["Rol, görev, hesap durumu ve yönetici hiyerarşisi",
       "Takım lideri / operasyon müdürü / idari müdür alanları",
       "Görev grubu ataması yetkilendirmenin temeli",
       "Excel'e aktarma"],
      GREEN, "14 · KULLANICILAR")

ekran("Batch Girişi",
      "Toplu cihaz kabulü — Excel veya manuel, 1.817 satırlık ekran",
      "14_batchentry.png",
      ["Excel sütun eşleme modali ile esnek içe aktarma",
       "IMEI okutunca PhoneCheck'ten cihaz bilgisi otomatik doluyor",
       "Flow (Refurbish / Repair / Return / Battery) seçimi",
       "Yeni cihazlar 100 statüsünde açılıyor"],
      CYAN, "15 · BATCH GİRİŞİ")

ekran("Raporlar",
      "Stok, hareket ve üretim raporlarının merkezi",
      "15_raporlar.png",
      ["Tarih aralığı ve lokasyon filtreleri",
       "Excel çıktısı",
       "Malzeme tüketim ve üretim raporları",
       "Kritik stok listesi"],
      AMBER, "16 · RAPORLAR")

ekran("Schema Mapper",
      "Veritabanı şemasının etkileşimli görselleştirmesi",
      "16_schemamapper.png",
      ["70 tablo ve ilişkileri kanvas üzerinde",
       "Bezier bağlantı çizgileriyle yabancı anahtarlar",
       "Pan / zoom desteği (useCanvasPanZoom)",
       "Geliştirici aracı — şema keşfi için"],
      VIOLET, "17 · SCHEMA MAPPER")

# Tema karşılaştırma
s = new()
glow(s, 12.6, 0.4, 2.4, CYAN, 4)
chip(s, 0.75, 0.55, "18 · TEMA", CYAN)
txt(s, 0.75, 0.98, 11, 0.5, "Açık ve Koyu Tema", 25, WHITE, True)
txt(s, 0.75, 1.52, 11, 0.4,
    "Tek CSS dosyasından yönetilen çift tema — kullanıcı bazında hatırlanıyor", 11, MUTED)
# Ayni ekranin iki temadaki hali - farkli sayfalar konursa karsilastirma anlamini yitirir.
pic(s, shot("17_dashboard_acik.png"), 0.75, 2.05, 5.8)
pic(s, shot("02_dashboard.png"), 6.8, 2.05, 5.8)
txt(s, 0.75, 5.82, 5.8, 0.3, "AÇIK TEMA", 10, RGBColor(0xA9, 0xB2, 0xE3), True, spacing=1.6)
txt(s, 6.8, 5.82, 5.8, 0.3, "KOYU TEMA (VARSAYILAN)", 10, CYAN, True, spacing=1.6)
card(s, 0.75, 6.2, 11.85, 0.78, fill=PANEL2, line=BORDER, accent=CYAN)
txt(s, 1.1, 6.36, 11.2, 0.5,
    "Tailwind'in renk rampaları index.css içinde ezilerek 1500'den fazla sınıf kullanımı tek "
    "merkezden temaya uyduruluyor.\nTema tercihi kullanıcı adına göre saklanıyor: theme_<username>.",
    10, MUTED, line_spacing=1.25)
footer(s, "Tema")

# ════════════════════════════════════════════════════════════════════
section("05", "Mühendislik", "Köprü mimarisi, performans ve kalite kararları", ROSE)

# QWebChannel
s = new()
title_block(s, "KÖPRÜ", "Frontend ve backend nasıl konuşuyor?",
            "194 @Slot metodu, JSON string sözleşmesi, tek erişim noktası", ROSE)

card(s, 0.75, 2.7, 7.4, 2.35, fill=PANEL2, line=BORDER, accent=CYAN)
txt(s, 1.1, 2.95, 6.6, 0.3, "AKIŞ", 10, CYAN, True, spacing=1.6)
adimlar = [
    ("1", "React bileşeni", "api.getParts() çağırır"),
    ("2", "api.js sarmalayıcı", "backend.get_parts(callback)"),
    ("3", "QWebChannel", "Qt slotunu çağırır"),
    ("4", "web_bridge.py", "JSON string döner → JSON.parse"),
]
for i, (n, t, d) in enumerate(adimlar):
    yy = 3.35 + i * 0.42
    rect(s, 1.1, yy + 0.02, 0.26, 0.26, fill=CYAN, alpha=20, shape=MSO_SHAPE.OVAL, line=CYAN, lw=0.75)
    txt(s, 1.1, yy + 0.05, 0.26, 0.22, n, 9, CYAN, True, align=PP_ALIGN.CENTER, font=FONT_MONO)
    txt(s, 1.5, yy, 2.6, 0.3, t, 11, WHITE, True)
    txt(s, 4.15, yy, 3.6, 0.3, d, 10, MUTED, font=FONT_MONO)

card(s, 0.75, 5.25, 7.4, 1.6, fill=PANEL, line=RED, accent=RED)
txt(s, 1.1, 5.48, 6.6, 0.3, "⚠  KRİTİK KURAL", 10, RED, True, spacing=1.4)
txt(s, 1.1, 5.82, 6.7, 0.9,
    "@Slot kaç tip bildirdiyse JS tarafı O KADAR argüman + callback göndermek zorunda.\n"
    "Eksik argümanda QWebChannel çağrıyı hiç yapmaz: hata çıkmaz, Promise sonsuza kadar\n"
    "bekler ve ekran \"İşleniyor...\" halinde donar.",
    10.5, MUTED, line_spacing=1.25)

card(s, 8.45, 2.7, 4.15, 4.15, fill=PANEL2, line=BORDER, accent=VIOLET)
txt(s, 8.8, 2.95, 3.5, 0.3, "SÖZLEŞME", 10, VIOLET, True, spacing=1.6)
bullets(s, 8.8, 3.38, 3.5, 3.3, [
    "Python her zaman JSON string döner",
    "Sonuç şekli: { success, message?, ...veri }",
    "Son argüman daima callback",
    "191 API sarmalayıcı → 190 farklı slot",
    "Bağlantı yoksa mock backend devreye girer, ekran beyaz kalmaz",
], size=10, gap=10, bullet_color=VIOLET)
footer(s, "QWebChannel Köprüsü")

# Performans
s = new()
title_block(s, "PERFORMANS", "Hız için alınan kararlar",
            "Ölçülebilir kazanç sağlayan dört optimizasyon", ROSE)
perf = [
    ("Rota bazlı kod bölme", "28 sayfa tek ~1.4 MB pakette birleşiyordu. Artık sadece "
     "ziyaret edilen ekran indiriliyor.", "lazy() + Suspense", CYAN),
    ("Vendor ayrıştırma", "React, router ve ikon kütüphanesi ayrı paketlere alındı; "
     "yeni derlemelerde tarayıcı bunları yeniden indirmiyor.", "manualChunks", GREEN),
    ("Önbellek başlıkları", "/assets hash'li dosyalar 1 yıl immutable, index.html asla "
     "önbelleğe alınmıyor.", "no-store + immutable", AMBER),
    ("Backend önbelleği", "Sık çağrılan get_users gibi slotlar 60 saniye önbellekleniyor, "
     "yazma işleminde anında geçersiz kılınıyor.", "TTL cache", VIOLET),
]
for i, (t, d, tag, c) in enumerate(perf):
    col, row = i % 2, i // 2
    x = 0.75 + col * 6.2
    y = 2.72 + row * 2.1
    card(s, x, y, 5.85, 1.9, accent=c)
    txt(s, x + 0.35, y + 0.28, 4.6, 0.35, t, 14, WHITE, True)
    txt(s, x + 0.35, y + 0.75, 5.15, 0.9, d, 10.5, MUTED, line_spacing=1.3)
    ww = 0.13 * len(tag) + 0.3
    rect(s, x + 0.35, y + 1.42, ww, 0.28, fill=c, alpha=16,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.5, line=c, lw=0.7)
    txt(s, x + 0.35, y + 1.475, ww, 0.22, tag, 8.5, c, True, align=PP_ALIGN.CENTER, font=FONT_MONO)
footer(s, "Performans")

# ════════════════════════════════════════════════════════════════════
section("06", "Kalite ve Yol Haritası", "Düzeltilen hatalar ve sıradaki adımlar", RED)

# Düzeltilen hatalar
s = new()
title_block(s, "KALİTE", "Son dönemde düzeltilen kritik hatalar",
            "Hepsi kaynak kodda doğrulandı ve yeniden derlendi", RED)
hatalar = [
    ("Statü akışı kilitlenmesi", "batch_entries.flow alanı 7.644 cihazda kısa adı "
     "(\"Refurbish\") tutuyordu; kural kodu (\"To refurbish\") ile karşılaştırıyordu. "
     "Neredeyse hiçbir onarım kapanamıyordu.", "7.644 cihaz etkilenmişti", RED),
    ("Test ekranları donuyordu", "submit_test_result 8 argüman bekliyor, api.js 7 "
     "gönderiyordu. Ara Test ve Son Test ekranları sessizce kilitleniyordu.",
     "2 ekran", AMBER),
    ("Servis kaydı kaydedilemiyordu", "create/update_service_record çağrılarında "
     "imei_number argümanı atlanmıştı.", "2 işlem", AMBER),
    ("İç sayfalar yenilenince açılmıyordu", "vite base './' idi; çok segmentli rotalarda "
     "(/onarim-havuzu/BATTERY) varlık yolları 404 veriyor ve React hiç açılmıyordu.",
     "11 rota", RED),
]
for i, (t, d, n, c) in enumerate(hatalar):
    y = 2.6 + i * 1.08
    card(s, 0.75, y, 11.85, 0.96, accent=c)
    rect(s, 1.1, y + 0.32, 0.32, 0.32, fill=c, alpha=20, shape=MSO_SHAPE.OVAL, line=c, lw=0.75)
    txt(s, 1.1, y + 0.36, 0.32, 0.26, "✓", 11, c, True, align=PP_ALIGN.CENTER)
    txt(s, 1.62, y + 0.14, 4.2, 0.32, t, 13, WHITE, True)
    txt(s, 1.62, y + 0.52, 8.3, 0.4, d, 9.5, MUTED, line_spacing=1.2)
    txt(s, 10.3, y + 0.34, 2.05, 0.3, n, 10, c, True, align=PP_ALIGN.RIGHT, font=FONT_MONO)
footer(s, "Düzeltilen Hatalar")

# Yol haritası
s = new()
title_block(s, "YOL HARİTASI", "Sıradaki adımlar",
            "Kod incelemesinde tespit edilen, öncelik sırasına göre", RED)
yol = [
    ("Yüksek", "Rota koruması", "MainLayout altındaki rotalarda oturum kontrolü yok; "
     "adres yazarak ekranlara erişilebiliyor.", RED),
    ("Yüksek", "Argüman uyum testi", "@Slot ↔ api.js argüman sayısı kontrolü CI adımına "
     "eklenmeli; bu hata sınıfı sessiz kilitlenme yaratıyor.", RED),
    ("Orta", "Ölü kod temizliği", "Table.jsx ve DbErrorModal.jsx hiç kullanılmıyor; "
     "axios, zustand, qwebchannel paketleri gereksiz.", AMBER),
    ("Orta", "NotificationToast ortaklaştırma", "7 dosyada birebir kopyalanmış durumda.", AMBER),
    ("Düşük", "WorkOrders bölünmesi", "3.426 satırlık dosya sekme başına ayrılabilir.", CYAN),
    ("Düşük", "Eski Pil eşlemesi", "PhoneCheck 10 cihazın hiçbirinde \"Old Battery\" "
     "döndürmüyor; alan eşlemesi netleşmeli.", CYAN),
]
for i, (onc, t, d, c) in enumerate(yol):
    y = 2.58 + i * 0.72
    card(s, 0.75, y, 11.85, 0.62, fill=PANEL if i % 2 == 0 else PANEL2, line=BORDER, accent=c, aw=0.04)
    rect(s, 1.05, y + 0.17, 0.8, 0.28, fill=c, alpha=16,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.5, line=c, lw=0.7)
    txt(s, 1.05, y + 0.222, 0.8, 0.22, onc, 8, c, True, align=PP_ALIGN.CENTER)
    txt(s, 2.05, y + 0.15, 3.3, 0.3, t, 12, WHITE, True)
    txt(s, 5.45, y + 0.175, 6.9, 0.4, d, 9.5, MUTED)
footer(s, "Yol Haritası")

# Kapanış
s = new(BG2)
grad(s, 0, 0, 13.333, 7.5, RGBColor(0x0B, 0x10, 0x22), BG, angle=45)
glow(s, 6.6, 3.6, 4.4, CYAN, 6)
if os.path.exists(logo):
    from PIL import Image as _I2
    lw2, lh2 = _I2.open(logo).size
    LW = 2.55
    s.shapes.add_picture(logo, Inches((13.333 - LW) / 2), Inches(1.32),
                         Inches(LW), Inches(LW * lh2 / lw2))
txt(s, 0.75, 4.05, 11.85, 0.8, "Teşekkürler", 44, WHITE, True, align=PP_ALIGN.CENTER)
txt(s, 0.75, 4.9, 11.85, 0.5,
    "RemaLab WMS  ·  Cihaz yaşam döngüsünün tamamı tek sistemde",
    14, MUTED, align=PP_ALIGN.CENTER)
rect(s, 6.17, 5.55, 1.0, 0.05, fill=CYAN)
txt(s, 0.75, 5.9, 11.85, 0.4,
    "Detaylı teknik döküman:  docs/FRONTEND.md   ·   docs/SISTEM_DOKUMANTASYONU.md",
    10.5, DIM, align=PP_ALIGN.CENTER, font=FONT_MONO)
footer(s, "Kapanış")

prs.save(OUT)
print("OK ->", OUT)
print("Slayt sayisi:", len(prs.slides.__iter__.__self__._sldIdLst))
